#!/usr/bin/env python3
"""Build the final poster and project deck from real product screenshots."""

from __future__ import annotations

import argparse
from pathlib import Path

from PIL import Image, ImageDraw, ImageEnhance, ImageFilter, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
SCREENSHOT_ROOT = ROOT / "submission/assets/screenshots"
DEFAULT_OUTPUT = ROOT / "submission/generated"
FONT_PATH = Path.home() / ".local/share/fonts/NotoSansCJKsc-Regular.otf"

BG = RGBColor(8, 13, 12)
PANEL = RGBColor(18, 25, 20)
PANEL_ALT = RGBColor(13, 30, 28)
INK = RGBColor(243, 239, 222)
MUTED = RGBColor(177, 183, 169)
GOLD = RGBColor(221, 164, 64)
TEAL = RGBColor(68, 177, 163)
RED = RGBColor(204, 95, 79)
SOFT_GOLD = RGBColor(60, 48, 26)
SOFT_TEAL = RGBColor(20, 49, 44)
SOFT_RED = RGBColor(54, 31, 28)


def font(size: int) -> ImageFont.FreeTypeFont:
    return ImageFont.truetype(str(FONT_PATH), size)


def cover_crop(image: Image.Image, size: tuple[int, int]) -> Image.Image:
    target_w, target_h = size
    ratio = max(target_w / image.width, target_h / image.height)
    resized = image.resize((round(image.width * ratio), round(image.height * ratio)))
    left = (resized.width - target_w) // 2
    top = (resized.height - target_h) // 2
    return resized.crop((left, top, left + target_w, top + target_h))


def draw_text(draw: ImageDraw.ImageDraw, xy: tuple[int, int], text: str, size: int, fill: str, spacing: int = 8) -> None:
    draw.multiline_text(xy, text, font=font(size), fill=fill, spacing=spacing)


def build_poster(output: Path) -> Path:
    canvas = cover_crop(Image.open(SCREENSHOT_ROOT / "main_battle.png").convert("RGB"), (1920, 1080))
    canvas = ImageEnhance.Contrast(canvas).enhance(1.08)
    canvas = ImageEnhance.Color(canvas).enhance(0.88)
    overlay = Image.new("RGBA", canvas.size, (0, 0, 0, 0))
    pixels = overlay.load()
    for x in range(1920):
        alpha = int(220 * max(0.0, 1.0 - x / 1320)) + 28
        for y in range(1080):
            pixels[x, y] = (4, 8, 7, min(alpha, 235))
    canvas = Image.alpha_composite(canvas.convert("RGBA"), overlay)
    draw = ImageDraw.Draw(canvas)
    draw.rounded_rectangle((108, 94, 548, 142), radius=10, fill=(24, 42, 34, 235), outline=(68, 177, 163, 220), width=2)
    draw_text(draw, (132, 102), "VIBE PLAYING · AI 原生游戏", 27, "#a9e8d8")
    draw_text(draw, (102, 184), "Compiler", 128, "#fff8df")
    draw_text(draw, (108, 350), "从 Vibe Coding\n到 Vibe Playing", 55, "#f4e9c8", 14)
    draw_text(draw, (112, 514), "自然语言意图  →  结构化解析  →  校验与模拟  →  Session 激活", 26, "#ced8cb")

    chips = [("防御塔", GOLD), ("陷阱", TEAL), ("支援道具", RED), ("三个编译世界", RGBColor(140, 125, 216))]
    x = 112
    for label, color in chips:
        width = 54 + len(label) * 31
        draw.rounded_rectangle((x, 604, x + width, 658), radius=10, fill=(10, 16, 14, 220), outline=tuple(color), width=2)
        draw_text(draw, (x + 24, 613), label, 25, "#f6f1df")
        x += width + 18

    draw.rounded_rectangle((112, 790, 1010, 968), radius=12, fill=(8, 13, 12, 222), outline=(196, 152, 70, 150), width=2)
    draw_text(draw, (146, 824), "不是让 AI 陪你聊游戏，而是让 AI 把意图变成玩法。", 28, "#e7c478")
    draw_text(draw, (146, 875), "Schema / 语义校验  /  确定性模拟  /  Promotion Gate  /  Session 隔离", 24, "#c5cec2")
    draw_text(draw, (146, 918), "可验证、可执行、可回退，才是 AI 原生游戏的底座。", 25, "#fff8df")
    draw_text(draw, (1630, 1016), "TEAM COMPILER", 23, "#dcb66b")
    output.parent.mkdir(parents=True, exist_ok=True)
    canvas.convert("RGB").save(output, quality=91, optimize=True, progressive=True)
    return output


def set_slide_bg(slide, color: RGBColor = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size: int, color: RGBColor = INK, bold: bool = False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    for run in paragraph.runs:
        run.font.name = "Noto Sans CJK SC"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_title(slide, eyebrow: str, title: str, subtitle: str = "") -> None:
    add_text(slide, eyebrow, 0.7, 0.42, 5.5, 0.3, 11, GOLD, True)
    add_text(slide, title, 0.7, 0.76, 11.9, 0.72, 30, INK, True)
    if subtitle:
        add_text(slide, subtitle, 0.72, 1.48, 11.6, 0.48, 14, MUTED)


def add_picture_cover(slide, path: Path, x: float, y: float, w: float, h: float) -> None:
    with Image.open(path) as image:
        image_ratio = image.width / image.height
    box_ratio = w / h
    picture = slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    if image_ratio > box_ratio:
        visible = box_ratio / image_ratio
        crop = (1 - visible) / 2
        picture.crop_left = crop
        picture.crop_right = crop
    else:
        visible = image_ratio / box_ratio
        crop = (1 - visible) / 2
        picture.crop_top = crop
        picture.crop_bottom = crop


def add_rect(slide, x: float, y: float, w: float, h: float, fill: RGBColor, line: RGBColor | None = None) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line


def add_rule(slide, x: float, y: float, w: float, color: RGBColor = GOLD, h: float = 0.035) -> None:
    add_rect(slide, x, y, w, h, color)


def add_badge(slide, text: str, x: float, y: float, w: float, color: RGBColor, fill: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.36))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = color
    shape.line.transparency = 32
    add_text(slide, text, x, y + 0.075, w, 0.2, 9, color, True, PP_ALIGN.CENTER)


def add_bullet(slide, text: str, x: float, y: float, w: float, color: RGBColor = MUTED, size: int = 13) -> None:
    add_text(slide, "•", x, y - 0.01, 0.24, 0.28, size + 2, GOLD, True)
    add_text(slide, text, x + 0.3, y, w - 0.3, 0.5, size, color)


def add_column_header(slide, kicker: str, title: str, x: float, y: float, w: float, accent: RGBColor, title_size: int = 22) -> None:
    add_rule(slide, x, y, w, accent)
    add_text(slide, kicker, x, y + 0.18, w, 0.24, 10, accent, True)
    add_text(slide, title, x, y + 0.54, w, 0.55, title_size, INK, True)


def add_image_frame(slide, path: Path, x: float, y: float, w: float, h: float, accent: RGBColor = MUTED) -> None:
    add_rect(slide, x - 0.025, y - 0.025, w + 0.05, h + 0.05, accent)
    add_picture_cover(slide, path, x, y, w, h)


def add_card(slide, x: float, y: float, w: float, h: float, title: str, body: str, accent: RGBColor = GOLD) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = accent; shape.line.transparency = 38
    add_text(slide, title, x + 0.25, y + 0.22, w - 0.5, 0.38, 16, INK, True)
    add_text(slide, body, x + 0.25, y + 0.72, w - 0.5, h - 0.88, 11, MUTED)


def add_footer(slide, index: int) -> None:
    add_text(slide, "Compiler · Tencent Cloud Hackathon", 0.72, 7.18, 5.5, 0.18, 8, MUTED)
    add_text(slide, f"{index:02d}", 12.1, 7.16, 0.5, 0.2, 8, GOLD, True, PP_ALIGN.RIGHT)


def build_deck(output: Path, poster: Path) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank); set_slide_bg(slide)
    add_picture_cover(slide, poster, 0, 0, 13.333, 7.5)

    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, "01 · AI 原生游戏范式", "Vibe Coding 之后，是 Vibe Playing", "把“自然语言 → 可运行软件”的范式，迁移为“玩家意图 → 可验证玩法 → 当前 Session”。")
    add_column_header(slide, "VIBE CODING", "意图变成软件", 0.75, 2.15, 3.25, TEAL)
    add_text(slide, "自然语言", 0.75, 3.15, 3.25, 0.35, 15, INK, True)
    add_text(slide, "→ 结构化产物 → 可运行应用", 0.75, 3.62, 3.25, 0.4, 13, MUTED)
    add_text(slide, "人描述目标，AI 与工程系统完成实现。", 0.75, 4.35, 3.25, 0.65, 12, MUTED)

    add_column_header(slide, "PARADIGM SHIFT", "把生成接入运行时", 5.02, 2.15, 3.25, GOLD)
    add_text(slide, "理解只是起点", 5.02, 3.15, 3.25, 0.35, 15, INK, True)
    add_text(slide, "→ 校验 → 模拟 → 晋升激活", 5.02, 3.62, 3.25, 0.4, 13, MUTED)
    add_text(slide, "AI 输出必须穿过游戏规则与质量门。", 5.02, 4.35, 3.25, 0.65, 12, MUTED)

    add_column_header(slide, "VIBE PLAYING", "意图变成玩法", 9.3, 2.15, 3.25, RED)
    add_text(slide, "自然语言意图", 9.3, 3.15, 3.25, 0.35, 15, INK, True)
    add_text(slide, "→ Runtime asset → 当前 Session", 9.3, 3.62, 3.25, 0.4, 13, MUTED)
    add_text(slide, "玩家提出想法，并在战场上观察结果。", 9.3, 4.35, 3.25, 0.65, 12, MUTED)

    add_rect(slide, 0.75, 5.45, 11.8, 1.03, PANEL_ALT)
    add_badge(slide, "本作验证命题", 1.0, 5.76, 1.45, TEAL, SOFT_TEAL)
    add_text(slide, "AI 原生游戏的关键，不是生成更多文案，而是改变玩家与系统的交互合同。", 2.72, 5.72, 9.45, 0.42, 17, INK, True)
    add_footer(slide, 2)

    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, "02 · 玩家体验", "玩家眼中的一条完整链路", "技术管线留在后台；前台只有世界、工坊，以及新能力进入战场的结果。")
    add_badge(slide, "已实现", 11.25, 0.45, 1.25, TEAL, SOFT_TEAL)
    add_image_frame(slide, SCREENSHOT_ROOT / "main_map.png", 0.7, 2.08, 3.72, 3.72, GOLD)
    add_image_frame(slide, SCREENSHOT_ROOT / "main_workshop.png", 4.81, 2.08, 3.72, 3.72, TEAL)
    add_image_frame(slide, SCREENSHOT_ROOT / "main_battle.png", 8.92, 2.08, 3.72, 3.72, RED)
    add_text(slide, "→", 4.45, 3.63, 0.32, 0.4, 20, MUTED, True, PP_ALIGN.CENTER)
    add_text(slide, "→", 8.56, 3.63, 0.32, 0.4, 20, MUTED, True, PP_ALIGN.CENTER)
    add_text(slide, "01 进入世界与战略地图", 0.82, 5.96, 3.48, 0.32, 12, GOLD, True, PP_ALIGN.CENTER)
    add_text(slide, "02 用自然语言提出意图", 4.93, 5.96, 3.48, 0.32, 12, TEAL, True, PP_ALIGN.CENTER)
    add_text(slide, "03 激活样品并观察行为", 9.04, 5.96, 3.48, 0.32, 12, RED, True, PP_ALIGN.CENTER)
    add_text(slide, "三个编译世界均可进入这条完整页面链路。", 2.15, 6.54, 9.05, 0.35, 15, INK, True, PP_ALIGN.CENTER)
    add_footer(slide, 3)

    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, "03 · AI 编译器", "从一句话，到一个可激活的 Session", "自然语言输入可以自由；进入运行时之前，每一步都必须结构化、可验证、可审计。")
    add_badge(slide, "已实现", 11.25, 0.45, 1.25, TEAL, SOFT_TEAL)
    steps = [
        ("01", "Intent", "玩家自然语言\n意图"),
        ("02", "Parse", "Provider 输出\n结构化候选"),
        ("03", "Validate", "Schema + 语义\n约束校验"),
        ("04", "Simulate", "确定性模拟\n与预算检查"),
        ("05", "Promote", "晋升报告\n决定可否发布"),
        ("06", "Activate", "Runtime asset\n进入 Session"),
    ]
    for i, (number, name, body) in enumerate(steps):
        x = 0.65 + i * 2.08
        add_rect(slide, x, 2.28, 1.76, 2.35, PANEL, [TEAL, GOLD, RED][i % 3])
        add_rule(slide, x, 2.28, 1.76, [TEAL, GOLD, RED][i % 3], 0.08)
        add_text(slide, number, x + 0.18, 2.57, 0.45, 0.25, 10, [TEAL, GOLD, RED][i % 3], True)
        add_text(slide, name, x + 0.18, 2.96, 1.4, 0.35, 16, INK, True)
        add_text(slide, body, x + 0.18, 3.55, 1.4, 0.7, 10, MUTED)
        if i < len(steps) - 1:
            add_text(slide, "→", x + 1.78, 3.23, 0.28, 0.4, 17, MUTED, True, PP_ALIGN.CENTER)
    add_rule(slide, 0.7, 5.12, 5.7, TEAL)
    add_text(slide, "模型边界", 0.7, 5.34, 1.3, 0.3, 12, TEAL, True)
    add_text(slide, "Provider 只提交候选，不直接改写世界状态。", 0.7, 5.82, 5.7, 0.55, 14, INK, True)
    add_rule(slide, 6.85, 5.12, 5.75, GOLD)
    add_text(slide, "发布证据", 6.85, 5.34, 1.3, 0.3, 12, GOLD, True)
    add_text(slide, "PromotionReport + ActivationReceipt 记录晋升与激活。", 6.85, 5.82, 5.75, 0.55, 14, INK, True)
    add_footer(slide, 4)

    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, "04 · 真实运行闭环", "不是生成一段描述，而是改写一次可观察的战斗", "塔、陷阱与支援道具都完成了从候选到战斗行为的闭环。")
    add_badge(slide, "已实现", 11.25, 0.45, 1.25, TEAL, SOFT_TEAL)
    add_text(slide, "对象", 0.9, 2.16, 1.2, 0.3, 10, GOLD, True)
    add_text(slide, "运行时行为", 2.35, 2.16, 2.8, 0.3, 10, GOLD, True)
    add_text(slide, "可验证证据", 5.35, 2.16, 2.2, 0.3, 10, GOLD, True)
    rows = [
        ("防御塔", "部署 → 索敌 → 伤害 / 连锁 → 结算", "费用、射程、目标数进入行为 ABI", TEAL),
        ("陷阱", "部署 → 触发 → 范围伤害 / 减速 → 失效", "触发窗口、作用半径、生命周期可模拟", GOLD),
        ("支援道具", "选点 → 释放 → 范围效果 → 冷却", "落点、费用、范围与冷却进入运行时", RED),
    ]
    for i, (name, behavior, evidence, accent) in enumerate(rows):
        y = 2.58 + i * 1.12
        add_rect(slide, 0.72, y, 6.93, 0.86, PANEL if i % 2 == 0 else PANEL_ALT)
        add_rule(slide, 0.72, y, 0.07, accent, 0.86)
        add_text(slide, name, 0.94, y + 0.25, 1.15, 0.3, 14, accent, True)
        add_text(slide, behavior, 2.35, y + 0.25, 2.75, 0.35, 12, INK, True)
        add_text(slide, evidence, 5.35, y + 0.18, 2.05, 0.52, 10, MUTED)
    add_image_frame(slide, SCREENSHOT_ROOT / "main_battle.png", 8.05, 2.12, 4.55, 4.1, RED)
    add_badge(slide, "RUNTIME BATTLE", 8.28, 5.62, 1.65, RED, SOFT_RED)
    add_text(slide, "候选只有穿过校验、模拟与晋升，才会成为当前 Session 的玩法资产。", 0.8, 6.33, 7.0, 0.55, 13, INK, True)
    add_footer(slide, 5)

    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, "05 · Provider 与安全降级", "多 Provider，不把不确定性传给玩家", "模型可以替换、失败或返回不合规候选；统一质量门与安全 fallback 保持玩法可用。")
    add_badge(slide, "已实现", 11.25, 0.45, 1.25, TEAL, SOFT_TEAL)
    add_rect(slide, 0.75, 2.08, 11.82, 1.45, PANEL_ALT)
    flow = [("玩家意图", TEAL), ("多 Provider 路由", GOLD), ("结构化候选", RED), ("统一质量门", TEAL), ("Session / fallback", GOLD)]
    for i, (label, accent) in enumerate(flow):
        x = 1.0 + i * 2.32
        add_text(slide, f"0{i + 1}", x, 2.42, 0.38, 0.25, 9, accent, True)
        add_text(slide, label, x, 2.82, 1.85, 0.35, 13, INK, True)
        if i < len(flow) - 1:
            add_text(slide, "→", x + 1.88, 2.76, 0.3, 0.35, 15, MUTED, True, PP_ALIGN.CENTER)
    safeguards = [
        ("输出边界", "Schema / 语义校验 / 白名单", TEAL),
        ("执行边界", "确定性模拟 / 预算 / 晋升门", GOLD),
        ("状态边界", "Session 隔离 / 激活回执 / 可回退", RED),
        ("降级边界", "候选不可用 → 安全 fallback", TEAL),
    ]
    for i, (title, body, accent) in enumerate(safeguards):
        x = 0.75 + i * 3.03
        add_rule(slide, x, 4.2, 2.7, accent)
        add_text(slide, title, x, 4.5, 2.7, 0.35, 15, accent, True)
        add_text(slide, body, x, 5.12, 2.7, 0.8, 11, MUTED)
    add_text(slide, "Provider 负责提出可能性；Compiler 负责决定什么能进入游戏。", 1.4, 6.25, 10.55, 0.5, 16, INK, True, PP_ALIGN.CENTER)
    add_footer(slide, 6)

    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, "06 · 地图编译", "地图结构闭环，视觉仍用 reviewed fallback", "逻辑事实与视觉表现分离：路线、塔位、出生点和目标不由一张生成图反推。")
    add_badge(slide, "结构管线 · 已实现", 8.85, 0.45, 1.8, TEAL, SOFT_TEAL)
    add_badge(slide, "当前视觉 · FALLBACK", 10.82, 0.45, 1.78, GOLD, SOFT_GOLD)
    add_image_frame(slide, SCREENSHOT_ROOT / "main_map.png", 0.7, 2.08, 5.55, 4.48, TEAL)
    add_column_header(slide, "STRUCTURE / 已实现", "地图分层编译闭环", 6.72, 2.08, 5.88, TEAL)
    map_stages = [
        ("01", "LogicMapIR", "拓扑、路径、部署槽位、出生点与目标可测试"),
        ("02", "VisualMaterialPack", "材质候选与逻辑坐标解耦，保留审查入口"),
        ("03", "MapVisualRuntimePackage", "确定性合成、运行时高亮与 fallback 契约"),
    ]
    for i, (number, title, body) in enumerate(map_stages):
        y = 3.22 + i * 0.92
        add_text(slide, number, 6.72, y, 0.45, 0.28, 10, [TEAL, GOLD, RED][i], True)
        add_text(slide, title, 7.35, y - 0.02, 2.2, 0.32, 14, INK, True)
        add_text(slide, body, 9.65, y - 0.02, 2.95, 0.55, 10, MUTED)
        if i < 2:
            add_rule(slide, 7.35, y + 0.58, 5.25, PANEL)
    add_rect(slide, 6.72, 5.92, 5.88, 0.64, SOFT_GOLD)
    add_text(slide, "边界说明", 6.94, 6.1, 1.05, 0.25, 10, GOLD, True)
    add_text(slide, "当前可玩地图视觉仍使用 reviewed fallback，不把候选视觉表述为已发布成果。", 8.22, 6.07, 4.1, 0.38, 10, INK, True)
    add_footer(slide, 7)

    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, "07 · 三个编译世界", "同一套 Compiler，进入三条完整页面链路", "仙侠、西幻与科幻拥有各自的世界书、开场、战略地图、工坊上下文与首关。")
    add_badge(slide, "已实现 · 3 WORLDS", 10.85, 0.45, 1.7, TEAL, SOFT_TEAL)
    world_maps = [
        ("xianxia_map.png", "云海断峰关 · 仙侠", "灵脉为路，符阵为城", TEAL),
        ("western_map.png", "石风边境领 · 西幻", "沼泽荒原上的最后堡垒", GOLD),
        ("scifi_map.png", "星锚轨道站 · 科幻", "轨道设施的模块化防线", RED),
    ]
    for index, (image, title, subtitle, accent) in enumerate(world_maps):
        x = 0.7 + index * 4.12
        add_image_frame(slide, SCREENSHOT_ROOT / image, x, 2.08, 3.7, 3.7, accent)
        add_rule(slide, x, 5.98, 3.7, accent)
        add_text(slide, title, x, 6.12, 3.7, 0.3, 13, accent, True, PP_ALIGN.CENTER)
        add_text(slide, subtitle, x, 6.5, 3.7, 0.25, 9, MUTED, False, PP_ALIGN.CENTER)
    add_text(slide, "共同链路：世界书 → 开场 → 战略地图 → 工坊 → 首关战斗", 3.0, 6.82, 7.35, 0.26, 10, INK, True, PP_ALIGN.CENTER)
    add_footer(slide, 8)

    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, "08 · AI 协作", "工具参与生产，质量门决定发布", "代码、结构化玩法候选和视觉候选来自不同工具；任何模型输出都不直接成为交付结果。")
    tool_columns = [
        ("CodeBuddy", "实现、重构\n与代码审查", TEAL),
        ("Codex / OpenCode", "架构、编排\n与集成验收", GOLD),
        ("DeepSeek / GLM", "结构化玩法候选\n与 Provider 协作", RED),
        ("Agnes / OpenAI", "图像模型：视觉候选\n与风格探索", TEAL),
    ]
    for i, (name, body, accent) in enumerate(tool_columns):
        x = 0.75 + i * 3.03
        add_column_header(slide, f"TOOL 0{i + 1}", name, x, 2.08, 2.7, accent, 18)
        add_text(slide, body, x, 3.4, 2.7, 0.75, 13, MUTED)
    add_rect(slide, 0.75, 4.75, 11.82, 1.48, PANEL_ALT)
    add_text(slide, "QUALITY GATE", 1.02, 5.05, 1.45, 0.25, 10, GOLD, True)
    add_text(slide, "Schema / 语义", 2.55, 5.45, 1.55, 0.3, 12, INK, True)
    add_text(slide, "→", 4.1, 5.42, 0.35, 0.3, 15, MUTED, True, PP_ALIGN.CENTER)
    add_text(slide, "模拟 / 晋升", 4.5, 5.45, 1.55, 0.3, 12, INK, True)
    add_text(slide, "→", 6.05, 5.42, 0.35, 0.3, 15, MUTED, True, PP_ALIGN.CENTER)
    add_text(slide, "多模态审查", 6.45, 5.45, 1.55, 0.3, 12, INK, True)
    add_text(slide, "→", 8.0, 5.42, 0.35, 0.3, 15, MUTED, True, PP_ALIGN.CENTER)
    add_text(slide, "自动测试", 8.4, 5.45, 1.35, 0.3, 12, INK, True)
    add_text(slide, "→", 9.75, 5.42, 0.35, 0.3, 15, MUTED, True, PP_ALIGN.CENTER)
    add_text(slide, "人工终审", 10.15, 5.45, 1.45, 0.3, 12, INK, True)
    add_text(slide, "统一验收对象：代码、Runtime package、发布证据", 1.0, 6.5, 11.35, 0.32, 12, MUTED, False, PP_ALIGN.CENTER)
    add_footer(slide, 9)

    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, "09 · NOW / NEXT", "已经发生的，与接下来要发生的", "本次 Demo 只把经过运行链路验证的能力标为“已实现”；长期世界演化明确属于后续规划。")
    add_rect(slide, 0.7, 2.05, 5.85, 4.73, SOFT_TEAL)
    add_badge(slide, "已实现 · NOW", 0.98, 2.35, 1.4, TEAL, PANEL_ALT)
    add_text(slide, "当前可验证能力", 0.98, 2.92, 4.95, 0.48, 22, INK, True)
    implemented = [
        "自然语言意图 → 结构化解析 → 校验 → 模拟 → 晋升激活 → Session",
        "塔 / 陷阱 / 支援道具均有可观察的行为闭环",
        "仙侠 / 西幻 / 科幻均可进入完整页面链路",
        "多 Provider 路由与安全 fallback",
        "地图分层编译已有结构闭环；当前视觉为 reviewed fallback",
    ]
    for i, item in enumerate(implemented):
        add_bullet(slide, item, 0.98, 3.63 + i * 0.58, 5.18, INK, 11)

    add_rect(slide, 6.78, 2.05, 5.85, 4.73, SOFT_GOLD)
    add_badge(slide, "后续规划 · NEXT", 7.06, 2.35, 1.65, GOLD, PANEL)
    add_text(slide, "不作为本次 Demo 成果", 7.06, 2.92, 4.95, 0.48, 22, INK, True)
    planned = [
        "更深剧情：长期角色关系与分支叙事",
        "任务系统：可编译目标、奖励与状态变化",
        "随机事件：在预算与世界规则内生成局势",
        "世界持续生长：跨局演化、记忆与长期进度",
    ]
    for i, item in enumerate(planned):
        add_bullet(slide, item, 7.06, 3.63 + i * 0.7, 5.18, INK, 12)
    add_rect(slide, 7.06, 6.25, 5.05, 0.36, PANEL)
    add_text(slide, "路线图，尚未标记为已实现。", 7.25, 6.33, 4.65, 0.2, 9, GOLD, True, PP_ALIGN.CENTER)
    add_footer(slide, 10)

    slide = prs.slides.add_slide(blank); set_slide_bg(slide)
    add_rect(slide, 7.08, 0, 6.253, 7.5, PANEL)
    add_picture_cover(slide, SCREENSHOT_ROOT / "main_battle.png", 7.1, 0, 6.233, 7.5)
    add_badge(slide, "VIBE PLAYING", 0.82, 0.94, 1.5, TEAL, SOFT_TEAL)
    add_text(slide, "从 Vibe Coding\n到 Vibe Playing", 0.8, 1.6, 5.85, 1.55, 37, INK, True)
    add_text(slide, "AI 原生游戏不只生成内容，\n它把玩家意图变为可验证、可执行、可回退的玩法。", 0.82, 3.58, 5.85, 1.25, 17, GOLD, True)
    add_rule(slide, 0.82, 5.3, 5.35, MUTED)
    add_text(slide, "结构化解析", 0.82, 5.62, 1.55, 0.3, 11, MUTED, True)
    add_text(slide, "校验与模拟", 2.65, 5.62, 1.55, 0.3, 11, MUTED, True)
    add_text(slide, "Session 激活", 4.5, 5.62, 1.55, 0.3, 11, MUTED, True)
    add_text(slide, "TEAM COMPILER", 0.82, 6.47, 3.0, 0.35, 13, TEAL, True)
    add_footer(slide, 11)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    return output


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--output-dir", type=Path, default=DEFAULT_OUTPUT)
    args = parser.parse_args()
    poster = build_poster(args.output_dir / "compiler_project_poster_16x9.jpg")
    deck = build_deck(args.output_dir / "Compiler-Project-Deck.pptx", poster)
    print(poster)
    print(deck)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
